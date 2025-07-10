import os
from dotenv import load_dotenv
load_dotenv()

import fitz  
from docx import Document as DocxDocument
from pptx import Presentation
from langchain_community.vectorstores import FAISS
from langchain_openai import AzureOpenAIEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document

from autogen import AssistantAgent, register_function


AZURE_OPENAI_API_KEY = os.getenv('AZURE_OPENAI_API_KEY')
AZURE_OPENAI_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_API_VERSION = os.getenv("AZURE_OPENAI_API_VERSION")
AZURE_OPENAI_EMBEDDING_DEPLOYMENT = os.getenv("AZURE_OPENAI_EMBEDDING_DEPLOYMENT")

openai_api_key = os.getenv("OPENAI_API_KEY")
openai_model = os.getenv("OPENAI_MODEL")


doc_assistant = AssistantAgent(
    name="DocAssistant",
    llm_config={
        "config_list": [
            {
                "model": openai_model,
                "api_key": openai_api_key
            }
        ]
    }
)


def extract_text_from_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    text = ""

    if ext == ".pdf":
        doc = fitz.open(file_path)
        for page in doc:
            text += page.get_text()

    elif ext == ".docx":
        doc = DocxDocument(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"

    elif ext == ".pptx":
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    else:
        raise ValueError(f"Unsupported file type: {ext}")

    return text

def build_index_from_file(file_path: str, persist_dir: str = "./rag_faiss_store"):
    full_text = extract_text_from_file(file_path)

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=100)
    documents = text_splitter.split_documents([Document(page_content=full_text)])

    embeddings = AzureOpenAIEmbeddings(
        azure_endpoint=AZURE_OPENAI_ENDPOINT,
        api_key=AZURE_OPENAI_API_KEY,
        api_version=AZURE_OPENAI_API_VERSION,
        azure_deployment=AZURE_OPENAI_EMBEDDING_DEPLOYMENT,
    )

    db = FAISS.from_documents(documents, embeddings)
    os.makedirs(persist_dir, exist_ok=True)
    db.save_local(persist_dir)


def retrieve_doc_context(query: str):
    embeddings = AzureOpenAIEmbeddings(
        azure_endpoint=AZURE_OPENAI_ENDPOINT,
        api_key=AZURE_OPENAI_API_KEY,
        api_version=AZURE_OPENAI_API_VERSION,
        azure_deployment=AZURE_OPENAI_EMBEDDING_DEPLOYMENT,
    )

    db = FAISS.load_local("rag_faiss_store", embeddings, allow_dangerous_deserialization=True)
    docs = db.similarity_search(query, k=3)

    print("Similar Documents Found:")
    for doc in docs:
        print(doc)

    return "\n\n".join([doc.page_content for doc in docs]), docs


register_function(
    retrieve_doc_context,
    caller=doc_assistant,
    executor=doc_assistant,
    description="Retrieve relevant context from the indexed documents (PDF, DOCX, PPTX) based on the user's query."
)


if __name__ == "__main__":
    
    build_index_from_file("./docs/NehaResume.docx")
    
    
    query = "What is the education of Neha?"
    context, sources = retrieve_doc_context(query)
    print("Retrieved Context:\n", context)
